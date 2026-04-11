from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from unifoli_render.formats.base import BaseRenderer
from unifoli_render.markdown import markdown_lines_to_bullets, split_markdown_sections
from unifoli_render.models import RenderArtifact, RenderBuildContext
from unifoli_render.template_registry import build_provenance_appendix_lines
from unifoli_shared.paths import to_stored_path


def _hex_to_rgb(value: str) -> RGBColor:
    cleaned = value.lstrip("#")
    return RGBColor(int(cleaned[0:2], 16), int(cleaned[2:4], 16), int(cleaned[4:6], 16))


class PptxRenderer(BaseRenderer):
    extension = ".pptx"
    implementation_level = "python-pptx"

    def render(self, context: RenderBuildContext, output_path: str | Path) -> RenderArtifact:
        self._build_presentation(context, output_path)

        relative_path = to_stored_path(output_path)
        return RenderArtifact(
            absolute_path=str(output_path),
            relative_path=relative_path,
            message="PPTX presentation created with python-pptx.",
        )

    def _build_presentation(self, context: RenderBuildContext, output_path) -> None:
        template = context.resolve_template()
        accent = _hex_to_rgb(template.preview.accent_color)

        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)

        title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        title_slide.shapes.title.text = context.project_title
        title_subtitle = title_slide.placeholders[1]
        title_subtitle.text = (
            f"{context.draft_title}\n"
            f"{template.label} template\n"
            f"Requested by: {context.requested_by or 'anonymous'}"
        )
        self._decorate_slide(title_slide, accent=accent, visual_priority=template.visual_priority)

        sections = split_markdown_sections(context.content_markdown)
        overview_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        overview_slide.shapes.title.text = template.preview.cover_title
        overview_body = overview_slide.placeholders[1].text_frame
        overview_body.clear()

        overview_points = [title for title, _ in sections if title][:6] or ["This draft is ready for expansion."]
        for index, bullet in enumerate(overview_points):
            paragraph = overview_body.paragraphs[0] if index == 0 else overview_body.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
            paragraph.font.size = Pt(24 if index == 0 else 20)
            paragraph.font.bold = index == 0

        self._decorate_slide(overview_slide, accent=accent, visual_priority=template.visual_priority)

        for section_title, lines in sections:
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = section_title or "Section"
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()
            bullets = markdown_lines_to_bullets(lines, max_items=7) or ["No content provided."]
            for index, bullet in enumerate(bullets):
                paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
                paragraph.text = bullet
                paragraph.level = 0
                paragraph.font.size = Pt(22 if index == 0 else 18)
                paragraph.font.name = "Aptos"
                paragraph.alignment = PP_ALIGN.LEFT
                if index == 0:
                    paragraph.font.bold = True

            self._decorate_slide(slide, accent=accent, visual_priority=template.visual_priority)

        appendix_lines = self._build_provenance_lines(context)
        if appendix_lines:
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "Provenance Appendix"
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()
            for index, line in enumerate(appendix_lines[:7]):
                paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
                paragraph.text = line
                paragraph.level = 0
                paragraph.font.size = Pt(18 if index == 0 else 14)
                if index == 0:
                    paragraph.font.bold = True
            self._decorate_slide(slide, accent=accent, visual_priority="low")

        presentation.save(str(output_path))

    @staticmethod
    def _build_provenance_lines(context: RenderBuildContext) -> list[str]:
        template = context.resolve_template()
        policy = context.export_policy
        if not policy.include_provenance_appendix or not template.supports_provenance_appendix:
            return []

        return [
            "Evidence basis used for this export",
            *build_provenance_appendix_lines(
                evidence_map=context.evidence_map,
                authenticity_log_lines=context.authenticity_log_lines,
                hide_internal=policy.hide_internal_provenance_on_final_export,
                max_evidence_items=4,
                max_authenticity_notes=2,
            ),
        ]

    @staticmethod
    def _decorate_slide(slide, *, accent: RGBColor, visual_priority: str) -> None:
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor(247, 248, 250)

        title = slide.shapes.title
        if title and title.text_frame and title.text_frame.paragraphs:
            first_run = title.text_frame.paragraphs[0].runs[0] if title.text_frame.paragraphs[0].runs else None
            if first_run:
                first_run.font.color.rgb = accent
                first_run.font.size = Pt(30 if visual_priority == "high" else 28)
                first_run.font.bold = True
